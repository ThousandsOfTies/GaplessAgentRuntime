#!/usr/bin/env python3
"""
5-slide storytelling deck. Personal, warm, not a tech spec.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette: clean light theme ─────────────────────────────────────────────
C_BG      = RGBColor(0xFB, 0xFB, 0xF9)   # off-white
C_PANEL   = RGBColor(0xF0, 0xF0, 0xEB)   # light warm gray
C_INK     = RGBColor(0x1A, 0x1A, 0x1A)   # near black
C_SUB     = RGBColor(0x55, 0x55, 0x55)   # gray
C_MUTED   = RGBColor(0x99, 0x99, 0x99)   # light gray
C_RED     = RGBColor(0xD0, 0x30, 0x30)   # accent red
C_BLUE    = RGBColor(0x1A, 0x6A, 0xC8)   # accent blue
C_GREEN   = RGBColor(0x1A, 0x9A, 0x60)   # accent green
C_ORANGE  = RGBColor(0xE0, 0x70, 0x20)   # accent orange
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

W = Inches(13.33)
H = Inches(7.5)


def prs_new():
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def tx(slide, text, l, t, w, h,
       size=18, bold=False, color=C_INK,
       align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name   = "Arial"
    return tb

def slide_num(slide, n):
    tx(slide, f"{n} / 5",
       W - Inches(1.2), H - Inches(0.4), Inches(1.0), Inches(0.3),
       size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)

def eyebrow(slide, text, color=C_MUTED):
    tx(slide, text, Inches(0.7), Inches(0.55),
       Inches(11), Inches(0.4), size=12, color=color, bold=True)

def left_bar(slide, color):
    rect(slide, Inches(0.45), Inches(0.45), Inches(0.08), H - Inches(0.9), color)


# ── Slide 1: あの時、見えていた ────────────────────────────────────────────
def s1(prs):
    slide = blank(prs)
    bg(slide)
    left_bar(slide, C_RED)
    slide_num(slide, 1)

    eyebrow(slide, "原点")

    tx(slide, "あの時、見えていた。",
       Inches(0.7), Inches(1.1), Inches(11), Inches(1.1),
       size=48, bold=True, color=C_INK)

    tx(slide,
       "放送機器のIP化が進んだとき、わかっていたことがあった。\n"
       "「専用ハードがコモディティ化する。日本メーカーの取り分が消える。」",
       Inches(0.7), Inches(2.4), Inches(10.5), Inches(1.3),
       size=20, color=C_SUB)

    # two columns: saw it / couldn't act
    rect(slide, Inches(0.7), Inches(3.9), Inches(5.3), Inches(2.7), C_PANEL)
    tx(slide, "見えていたこと",
       Inches(0.85), Inches(4.0), Inches(5.0), Inches(0.5),
       size=13, bold=True, color=C_RED)
    tx(slide,
       "・プロトコルを押さえた企業が残る\n"
       "・専用チップで変換を独占する\n"
       "・ソフトウェアが主戦場になる",
       Inches(0.85), Inches(4.55), Inches(5.0), Inches(1.9),
       size=15, color=C_INK)

    rect(slide, Inches(6.3), Inches(3.9), Inches(5.3), Inches(2.7), C_PANEL)
    tx(slide, "できなかったこと",
       Inches(6.45), Inches(4.0), Inches(5.0), Inches(0.5),
       size=13, bold=True, color=C_MUTED)
    tx(slide,
       "・予防線を張ること\n"
       "・儲かる仕組みを作ること\n"
       "・タイミングをつかむこと",
       Inches(6.45), Inches(4.55), Inches(5.0), Inches(1.9),
       size=15, color=C_SUB)

    tx(slide, "今度は違う動き方をする。",
       Inches(0.7), Inches(6.8), Inches(11), Inches(0.5),
       size=16, bold=True, italic=True, color=C_RED)


# ── Slide 2: 黒船は来ている ────────────────────────────────────────────────
def s2(prs):
    slide = blank(prs)
    bg(slide)
    left_bar(slide, C_ORANGE)
    slide_num(slide, 2)

    eyebrow(slide, "構造変化")

    tx(slide, "黒船は来ている。",
       Inches(0.7), Inches(1.1), Inches(11), Inches(1.1),
       size=48, bold=True, color=C_INK)

    tx(slide,
       "AIが「実装コスト」を汎用化する。\nコードを書く、ドライバを作る、テストを回す——全部が安くなる。",
       Inches(0.7), Inches(2.35), Inches(10.5), Inches(1.1),
       size=20, color=C_SUB)

    # dreamcast / iPad comparison
    rect(slide, Inches(0.7), Inches(3.6), Inches(11.7), Inches(0.08), C_PANEL)

    examples = [
        (C_ORANGE, "ドリームキャスト（1998）",
         "オンライン・携帯連携・独自メモリ\nすべて「正しいアイデア」だった\nタイミングと体力が合わなかった"),
        (C_BLUE,   "組み込み仮想化（2004〜）",
         "QEMUの時代から発想は完成していた\nコスト＞メリットで誰もやらなかった\n→ それが今、逆転しつつある"),
        (C_INK,    "AI × 組み込み（今）",
         "AIはクラウドにしか存在できない\n仮想HW環境がなければAIは\nドライバを書いても確認できない"),
    ]
    for i, (c, title, body) in enumerate(examples):
        x = Inches(0.7 + i * 4.0)
        y = Inches(3.8)
        rect(slide, x, y, Inches(0.08), Inches(3.0), c)
        tx(slide, title, x + Inches(0.2), y + Inches(0.05),
           Inches(3.6), Inches(0.55), size=13, bold=True, color=c)
        tx(slide, body, x + Inches(0.2), y + Inches(0.65),
           Inches(3.6), Inches(2.2), size=13, color=C_SUB)

    tx(slide, "正しいアイデアは、タイミングが来るまで待つしかない。",
       Inches(0.7), Inches(7.05), Inches(11), Inches(0.38),
       size=13, italic=True, color=C_MUTED)


# ── Slide 3: 今度こそタイミングが合った ───────────────────────────────────
def s3(prs):
    slide = blank(prs)
    bg(slide)
    left_bar(slide, C_BLUE)
    slide_num(slide, 3)

    eyebrow(slide, "なぜ今か")

    tx(slide, "今度こそ、タイミングが合った。",
       Inches(0.7), Inches(1.1), Inches(11), Inches(1.1),
       size=42, bold=True, color=C_INK)

    tx(slide, "2024〜25年、3つの変化が同時に起きた。",
       Inches(0.7), Inches(2.3), Inches(10), Inches(0.6),
       size=18, color=C_SUB)

    changes = [
        (C_BLUE,   "ISAの統一",
         "EC2 GravitonもRaspberry Pi 5も\n同じARM64になった",
         "エミュレーション不要。\nクロスコンパイルだけで\n本番と同じバイナリが動く",
         "根本課題が解消"),
        (C_GREEN,  "Cloud IDEの無料化",
         "GitHub Codespacesが\n無料枠で使えるようになった",
         "セットアップコストがゼロに。\n誰でも即日、\n組み込み開発を始められる",
         "参入障壁が消えた"),
        (C_ORANGE, "AIの登場",
         "Claude CodeがクラウドでC/Cを書く。\nでも実機には触れない",
         "仮想HW環境がなければ\nAIは組み込み開発に\n参加できない",
         "需要側の変化"),
    ]

    for i, (c, title, why, what, tag) in enumerate(changes):
        x = Inches(0.7 + i * 4.0)
        y = Inches(3.1)
        rect(slide, x, y, Inches(3.7), Inches(3.8), C_PANEL)
        rect(slide, x, y, Inches(3.7), Inches(0.08), c)

        tx(slide, title, x + Inches(0.2), y + Inches(0.2),
           Inches(3.3), Inches(0.5), size=15, bold=True, color=c)
        tx(slide, why, x + Inches(0.2), y + Inches(0.75),
           Inches(3.3), Inches(0.9), size=12, color=C_INK)
        tx(slide, what, x + Inches(0.2), y + Inches(1.7),
           Inches(3.3), Inches(1.3), size=12, color=C_SUB)

        rect(slide, x, y + Inches(3.35), Inches(3.7), Inches(0.45), c)
        tx(slide, tag, x + Inches(0.15), y + Inches(3.38),
           Inches(3.4), Inches(0.38), size=12, bold=True, color=C_WHITE)

    tx(slide,
       "ADSL → ネット普及。3G → スマホ普及。AI → 次の波。",
       Inches(0.7), Inches(7.08), Inches(11), Inches(0.35),
       size=12, italic=True, color=C_MUTED)


# ── Slide 4: もう動いている ────────────────────────────────────────────────
def s4(prs):
    slide = blank(prs)
    bg(slide)
    left_bar(slide, C_GREEN)
    slide_num(slide, 4)

    eyebrow(slide, "実証")

    tx(slide, "もう、動いている。",
       Inches(0.7), Inches(1.1), Inches(11), Inches(1.1),
       size=48, bold=True, color=C_INK)

    tx(slide,
       "ExperimentalDevEnv — 動くPoCがある。",
       Inches(0.7), Inches(2.35), Inches(10), Inches(0.55),
       size=18, color=C_SUB)

    # what's done
    rect(slide, Inches(0.7), Inches(3.1), Inches(6.8), Inches(3.6), C_PANEL)
    tx(slide, "できていること",
       Inches(0.85), Inches(3.2), Inches(6.4), Inches(0.45),
       size=13, bold=True, color=C_GREEN)
    done = [
        "GitHub Codespacesでクロスコンパイル（環境構築ゼロ）",
        "GPIO / I2C / SPIをクラウド上で仮想化",
        "ブラウザでハードウェアをリアルタイム操作",
        "同じバイナリをRaspberry Pi 5で動作確認済",
        "コード変更ゼロで実機展開",
    ]
    for j, d in enumerate(done):
        tx(slide, f"✓  {d}",
           Inches(0.85), Inches(3.75 + j * 0.55),
           Inches(6.4), Inches(0.5), size=13, color=C_INK)

    # key numbers
    nums = [
        (C_GREEN,  "0行",   "実機向け\nコード変更"),
        (C_BLUE,   "3種",   "HW I/F\n対応済"),
        (C_ORANGE, "$3/日", "EC2稼働\nコスト"),
    ]
    for i, (c, num, lbl) in enumerate(nums):
        x = Inches(7.9)
        y = Inches(3.1 + i * 1.25)
        rect(slide, x, y, Inches(4.3), Inches(1.1), C_PANEL)
        rect(slide, x, y, Inches(0.08), Inches(1.1), c)
        tx(slide, num, x + Inches(0.25), y + Inches(0.08),
           Inches(1.5), Inches(0.65), size=30, bold=True, color=c)
        tx(slide, lbl, x + Inches(1.85), y + Inches(0.2),
           Inches(2.2), Inches(0.7), size=12, color=C_SUB)

    tx(slide,
       "「作れる」の話ではなく、「作った」の話をしている。",
       Inches(0.7), Inches(7.05), Inches(11), Inches(0.38),
       size=14, bold=True, italic=True, color=C_GREEN)


# ── Slide 5: SONYがここで動く意味 ─────────────────────────────────────────
def s5(prs):
    slide = blank(prs)
    bg(slide)
    left_bar(slide, C_INK)
    slide_num(slide, 5)

    eyebrow(slide, "提案")

    tx(slide, "SONYがここで動く意味。",
       Inches(0.7), Inches(1.1), Inches(11), Inches(1.1),
       size=44, bold=True, color=C_INK)

    # why sony
    rect(slide, Inches(0.7), Inches(2.5), Inches(11.7), Inches(1.4), C_PANEL)
    tx(slide,
       "半導体・映像機器・ゲーム機・ロボティクス。\n"
       "すべてが組み込みの塊で、すべてがこの環境のユースケースに直撃する。",
       Inches(0.9), Inches(2.65), Inches(11.3), Inches(1.1),
       size=17, color=C_INK)

    # 3 steps
    steps = [
        (C_BLUE,  "Step 1  社内で使う",
         "まずこのチームで動かす\n実績とデータを積む\nリスクゼロ、コスト最小"),
        (C_ORANGE,"Step 2  社内展開",
         "他の開発チームへ広げる\n「SONYで使われている」\nが最大の証明になる"),
        (C_GREEN, "Step 3  外へ出る",
         "オープンコア＋コンサル\nデファクトスタンダードへ\n市場規模750億円"),
    ]
    for i, (c, title, body) in enumerate(steps):
        x = Inches(0.7 + i * 4.0)
        y = Inches(4.15)
        rect(slide, x, y, Inches(3.7), Inches(2.65), C_PANEL)
        rect(slide, x, y, Inches(3.7), Inches(0.08), c)
        tx(slide, title, x + Inches(0.2), y + Inches(0.2),
           Inches(3.3), Inches(0.5), size=14, bold=True, color=c)
        tx(slide, body, x + Inches(0.2), y + Inches(0.85),
           Inches(3.3), Inches(1.65), size=13, color=C_SUB)

        if i < 2:
            tx(slide, "→", Inches(4.28 + i * 4.0), y + Inches(1.2),
               Inches(0.5), Inches(0.4), size=20, color=C_MUTED,
               align=PP_ALIGN.CENTER)

    tx(slide,
       "あの時見えていたことを、今度は形にする。",
       Inches(0.7), Inches(7.05), Inches(11), Inches(0.38),
       size=15, bold=True, italic=True, color=C_INK)


# ── Main ──────────────────────────────────────────────────────────────────
def main():
    prs = prs_new()
    s1(prs)
    s2(prs)
    s3(prs)
    s4(prs)
    s5(prs)

    out = "/home/user/ExperimentalDevEnv/docs/ExperimentalDevEnv_Proposal.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
