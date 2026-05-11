#!/usr/bin/env python3
"""
Regenerate ExperimentalDevEnv intro presentation.
Narrative: "The idea isn't new — cost-benefit finally flipped."
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pptx.oxml.ns
from lxml import etree

# ── Palette ────────────────────────────────────────────────────────────────
C_BG      = RGBColor(0x0E, 0x14, 0x28)
C_PANEL   = RGBColor(0x18, 0x26, 0x42)
C_CYAN    = RGBColor(0x00, 0xC8, 0xFF)
C_PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
C_GREEN   = RGBColor(0x00, 0xD4, 0x9A)
C_ORANGE  = RGBColor(0xFF, 0x6B, 0x35)
C_YELLOW  = RGBColor(0xFF, 0xD0, 0x00)
C_RED     = RGBColor(0xFF, 0x45, 0x60)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED   = RGBColor(0x8A, 0xA4, 0xC8)

W = Inches(13.33)
H = Inches(7.5)


# ── Primitives ─────────────────────────────────────────────────────────────
def prs_new():
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color, line_color=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def tx(slide, text, l, t, w, h,
       size=18, bold=False, color=C_WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name  = "Arial"
    return tb

def tx_lines(slide, lines, l, t, w, h,
             sizes=None, colors=None, bolds=None,
             default_size=16, default_color=C_WHITE,
             align=PP_ALIGN.LEFT, spacing_after=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing_after:
            p.space_after = Pt(spacing_after)
        if not line:
            continue
        r = p.add_run()
        r.text  = line
        r.font.size  = Pt(sizes[i]  if sizes  and i < len(sizes)  else default_size)
        r.font.color.rgb = colors[i] if colors and i < len(colors) else default_color
        r.font.bold  = bolds[i] if bolds  and i < len(bolds)  else False
        r.font.name  = "Arial"
    return tb

def top_bar(slide, color=C_CYAN, height=Inches(0.07)):
    rect(slide, 0, 0, W, height, color)

def label(slide, text, color):
    """Small section label in top-left."""
    tx(slide, text, Inches(0.5), Inches(0.12), Inches(6), Inches(0.4),
       size=11, bold=True, color=color)

def heading(slide, text, sub=None):
    tx(slide, text, Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.85),
       size=32, bold=True, color=C_WHITE)
    if sub:
        tx(slide, sub, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.55),
           size=15, italic=True, color=C_MUTED)

def divider(slide, y, color=C_PANEL):
    rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.03), color)


# ── Slide 1: Title ─────────────────────────────────────────────────────────
def s_title(prs):
    slide = blank(prs)
    bg(slide)

    # left accent strip
    rect(slide, 0, 0, Inches(0.18), H, C_CYAN)
    # right panel
    rect(slide, Inches(7.2), 0, Inches(6.13), H, C_PANEL)

    # eyebrow
    tx(slide, "技術提案 ／ PoC",
       Inches(0.4), Inches(0.9), Inches(6.5), Inches(0.45),
       size=13, color=C_CYAN)

    # main title
    tx_lines(slide,
        ["クラウドネイティブ", "組み込み開発環境"],
        Inches(0.4), Inches(1.4), Inches(6.5), Inches(2.5),
        sizes=[44, 44], bolds=[True, True], colors=[C_WHITE, C_WHITE])

    # tagline
    tx(slide,
       "ハードウェアなしで開発を始め、\n同じバイナリを実機に展開する",
       Inches(0.4), Inches(3.85), Inches(6.5), Inches(1.2),
       size=18, italic=True, color=C_MUTED)

    # date / project
    tx(slide, "ExperimentalDevEnv  |  2026.05",
       Inches(0.4), Inches(6.8), Inches(6), Inches(0.45),
       size=12, color=C_MUTED)

    # right side: 3 steps
    steps = [
        (C_CYAN,   "1  Build",    "GitHub Codespaces\nクロスコンパイル"),
        (C_PURPLE, "2  Simulate", "AWS EC2 Graviton\nブラウザで HW 操作"),
        (C_GREEN,  "3  Deploy",   "Raspberry Pi 5\n同一バイナリで実機動作"),
    ]
    for i, (c, n, d) in enumerate(steps):
        y = Inches(1.2 + i * 1.9)
        rect(slide, Inches(7.5), y, Inches(5.5), Inches(1.55), C_BG)
        rect(slide, Inches(7.5), y, Inches(0.13), Inches(1.55), c)
        tx(slide, n, Inches(7.8), y + Inches(0.12),
           Inches(5.0), Inches(0.45), size=17, bold=True, color=c)
        tx(slide, d, Inches(7.8), y + Inches(0.6),
           Inches(5.0), Inches(0.7), size=14, color=C_MUTED)
        if i < 2:
            tx(slide, "↓", Inches(10.15), y + Inches(1.55),
               Inches(0.6), Inches(0.35), size=16, color=C_CYAN,
               align=PP_ALIGN.CENTER)


# ── Slide 2: "この発想、新しくない" ────────────────────────────────────────
def s_not_new(prs):
    slide = blank(prs)
    bg(slide)
    top_bar(slide, C_MUTED)

    label(slide, "01  前提の整理", C_MUTED)
    heading(slide, "この発想は、20年以上前からある",
            "QEMU・仮想化・エミュレーションはとっくに存在していた")

    # timeline
    tl = [
        ("2004", C_MUTED,   "QEMU 誕生",          "x86 エミュレーション。\nあらゆる CPU をソフトウェアで模倣可能に"),
        ("2010", C_MUTED,   "ARM Fast Models",    "ARM 公式の高速 CPU シミュレーター。\nチップより先にファームを開発できる"),
        ("2016", C_MUTED,   "Renode",             "IoT 向けマルチノード HW シミュレーター。\nボードまるごと仮想化"),
        ("2024", C_CYAN,    "本プロジェクト",      "なぜ今なのか？\n→ 次スライドで説明"),
    ]
    for i, (yr, c, title, body) in enumerate(tl):
        x = Inches(0.5 + i * 3.1)
        rect(slide, x, Inches(2.2), Inches(2.85), Inches(4.5), C_PANEL)
        rect(slide, x, Inches(2.2), Inches(2.85), Inches(0.1), c)
        tx(slide, yr,   x + Inches(0.15), Inches(2.35), Inches(2.5), Inches(0.45),
           size=22, bold=True, color=c)
        tx(slide, title, x + Inches(0.15), Inches(2.85), Inches(2.5), Inches(0.5),
           size=14, bold=True, color=C_WHITE)
        tx(slide, body,  x + Inches(0.15), Inches(3.4),  Inches(2.5), Inches(2.8),
           size=12, color=C_MUTED)

    # bottom emphasis
    rect(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.45), C_PANEL)
    tx(slide,
       "「仮想HWでテストする」は目新しいアイデアではない。  では、なぜ誰もやらなかったのか？",
       Inches(0.7), Inches(6.92), Inches(11.9), Inches(0.4),
       size=14, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)


# ── Slide 3: なぜやられなかったか ───────────────────────────────────────────
def s_why_not(prs):
    slide = blank(prs)
    bg(slide)
    top_bar(slide, C_RED)

    label(slide, "02  費用対効果の壁", C_RED)
    heading(slide, "コスト ＞ メリット  だったから",
            "3つの「割に合わない」が普及を阻んでいた")

    walls = [
        (C_RED,    "ISA の不一致",
         "クラウドは x86\n実機（RasPi等）は ARM\n\n"
         "→ エミュレーションが必要\n→ 動作が遅く信頼性も低い\n→ 「本番と同じ」にならない"),
        (C_ORANGE, "セットアップの複雑さ",
         "QEMU のビルド・設定に\n数日かかることも\n\n"
         "→ 導入コストが高すぎる\n→ 属人化・メンテ困難\n→ 結局「実機でいいや」"),
        (C_YELLOW, "クラウドコストの問題",
         "常時稼働のシミュレーター\nをクラウドで動かすのは\nコスト高だった\n\n"
         "→ 開発機 1 台の方が安い"),
    ]
    for i, (c, title, body) in enumerate(walls):
        x = Inches(0.5 + i * 4.1)
        rect(slide, x, Inches(2.1), Inches(3.8), Inches(4.8), C_PANEL)
        rect(slide, x, Inches(2.1), Inches(0.13), Inches(4.8), c)
        # big X
        tx(slide, "✕", x + Inches(2.9), Inches(2.2),
           Inches(0.7), Inches(0.7), size=28, bold=True, color=c,
           align=PP_ALIGN.CENTER)
        tx(slide, title, x + Inches(0.2), Inches(2.3),
           Inches(2.8), Inches(0.55), size=17, bold=True, color=c)
        tx(slide, body, x + Inches(0.2), Inches(2.95),
           Inches(3.4), Inches(3.8), size=13, color=C_MUTED)

    rect(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35), C_PANEL)
    tx(slide,
       "→  だからこそ、業界はずっと「実機があればそれでいい」を続けてきた",
       Inches(0.7), Inches(7.07), Inches(11.9), Inches(0.3),
       size=13, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER)


# ── Slide 4: なぜ今なのか（転換点）────────────────────────────────────────
def s_why_now(prs):
    slide = blank(prs)
    bg(slide)
    top_bar(slide, C_GREEN)

    label(slide, "03  転換点", C_GREEN)
    heading(slide, "2024〜2025年、3つの変化が同時に起きた",
            "この3年で「コスト ＜ メリット」の逆転が起きている")

    changes = [
        (C_CYAN,   "ISA の統一",
         "EC2 Graviton（AWS）が ARM64\nRaspberry Pi 5 も ARM64\n→ 同じ ISA\n\n"
         "エミュレーション不要。\nクロスコンパイルだけで\n本番と同じバイナリができる",
         "QEMU 時代の根本課題が解消"),
        (C_PURPLE, "Cloud IDE の無料化",
         "GitHub Codespaces が\n無料枠で利用可能に\n\n"
         "クロスコンパイル環境が\n即時・無料で手に入る\n→ セットアップコスト ≒ 0",
         "「数日かかる環境構築」が消滅"),
        (C_GREEN,  "AI の登場（最大の変化）",
         "Claude Code / Copilot が\nクラウド上でコードを書く\n\n"
         "AI は実機に触れない。\n→ クラウド上に HW 環境がないと\n   AI が開発できない",
         "「やりたい」から「やらないと」へ"),
    ]

    for i, (c, title, body, tag) in enumerate(changes):
        x = Inches(0.5 + i * 4.1)
        rect(slide, x, Inches(2.0), Inches(3.85), Inches(5.0), C_PANEL)
        rect(slide, x, Inches(2.0), Inches(3.85), Inches(0.12), c)
        # before → after
        tx(slide, "✕ → ✅", x + Inches(0.2), Inches(2.15),
           Inches(3.4), Inches(0.4), size=14, bold=True, color=c)
        tx(slide, title, x + Inches(0.2), Inches(2.6),
           Inches(3.4), Inches(0.55), size=16, bold=True, color=C_WHITE)
        tx(slide, body, x + Inches(0.2), Inches(3.2),
           Inches(3.4), Inches(3.0), size=12.5, color=C_MUTED)
        # tag
        rect(slide, x, Inches(6.75), Inches(3.85), Inches(0.5), c)
        tx(slide, tag, x + Inches(0.15), Inches(6.77),
           Inches(3.5), Inches(0.4), size=11.5, bold=True, color=C_BG)

    tx(slide,
       "★  この3つが揃ったのが今。本プロジェクトはそのタイミングを実証する PoC",
       Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
       size=14, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)


# ── Slide 5: ソリューション概要 ───────────────────────────────────────────
def s_solution(prs):
    slide = blank(prs)
    bg(slide)
    top_bar(slide, C_CYAN)

    label(slide, "04  ソリューション", C_CYAN)
    heading(slide, "ExperimentalDevEnv",
            "ハードなし → クラウドで検証 → 実機展開。3ステップのパイプライン PoC")

    steps = [
        (C_CYAN,   "☁  Build",
         "GitHub Codespaces",
         "aarch64-linux-gnu-gcc で ARM64 バイナリをビルド\n"
         "toolchain のセットアップはゼロ（devcontainer.json 1行）\n"
         "VSCode タスクで one-click ビルド",
         "コスト: 無料枠内"),
        (C_PURPLE, "⚡  Simulate",
         "AWS EC2 Graviton",
         "LD_PRELOAD shim → GPIO / SPI をインターセプト\n"
         "CUSE → /dev/i2c-1 仮想デバイスを生成\n"
         "Python ブリッジ → ブラウザで HW をリアルタイム操作",
         "コスト: t4g.micro ≈ $2〜3/日"),
        (C_GREEN,  "🔧  Deploy",
         "Raspberry Pi 5",
         "scp で同一バイナリを転送するだけ\n"
         "コード変更ゼロ・LD_PRELOAD 不要\n"
         "GPIO / I2C / SPI が実デバイスで動作",
         "コスト: デバイス不要（検証は EC2 で完了）"),
    ]

    for i, (c, icon_title, platform, body, cost) in enumerate(steps):
        x = Inches(0.5 + i * 4.1)
        rect(slide, x, Inches(2.0), Inches(3.85), Inches(5.3), C_PANEL)
        rect(slide, x, Inches(2.0), Inches(0.13), Inches(5.3), c)
        tx(slide, icon_title, x + Inches(0.2), Inches(2.1),
           Inches(3.5), Inches(0.5), size=18, bold=True, color=c)
        tx(slide, platform,  x + Inches(0.2), Inches(2.6),
           Inches(3.5), Inches(0.4), size=13, color=C_MUTED)
        tx(slide, body,      x + Inches(0.2), Inches(3.1),
           Inches(3.5), Inches(2.8), size=12.5, color=C_WHITE)
        rect(slide, x + Inches(0.2), Inches(6.9), Inches(3.4), Inches(0.28), C_BG)
        tx(slide, cost, x + Inches(0.25), Inches(6.91),
           Inches(3.2), Inches(0.25), size=11, color=c, bold=True)

        if i < 2:
            tx(slide, "→", Inches(4.28 + i * 4.1), Inches(4.4),
               Inches(0.5), Inches(0.5), size=22, bold=True, color=c,
               align=PP_ALIGN.CENTER)


# ── Slide 6: アーキテクチャ詳細 ───────────────────────────────────────────
def s_arch(prs):
    slide = blank(prs)
    bg(slide)
    top_bar(slide, C_PURPLE)

    label(slide, "05  仕組み", C_PURPLE)
    heading(slide, "アプリは何も知らない",
            "shim / CUSE が透過的にハードウェアを差し替える")

    # center: app box
    rect(slide, Inches(5.2), Inches(2.2), Inches(3.0), Inches(1.4), C_BG)
    rect(slide, Inches(5.2), Inches(2.2), Inches(3.0), Inches(0.1), C_WHITE)
    tx(slide, "アプリケーション\n(sensor_demo.c)",
       Inches(5.2), Inches(2.35), Inches(3.0), Inches(1.1),
       size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # left: shims
    shims = [
        (C_CYAN,   Inches(1.9), Inches(2.0), "GPIO Shim\ngpio_shim.so",
         "LD_PRELOAD\n/dev/gpiochip0 ioctl フック"),
        (C_PURPLE, Inches(1.9), Inches(3.8), "SPI Shim\nspi_shim.so",
         "LD_PRELOAD\n/dev/spidev0.0 ioctl フック"),
        (C_ORANGE, Inches(1.9), Inches(5.6), "I2C CUSE\ncuse_i2c",
         "FUSE3 で /dev/i2c-1\n仮想デバイス生成"),
    ]
    for c, x, y, t, d in shims:
        rect(slide, x, y, Inches(3.0), Inches(1.3), C_PANEL)
        rect(slide, x, y, Inches(0.1), Inches(1.3), c)
        tx(slide, t, x + Inches(0.15), y + Inches(0.08),
           Inches(2.7), Inches(0.5), size=13, bold=True, color=c)
        tx(slide, d, x + Inches(0.15), y + Inches(0.6),
           Inches(2.7), Inches(0.55), size=11, color=C_MUTED)
        # arrow to app
        tx(slide, "→", Inches(4.9), y + Inches(0.45),
           Inches(0.4), Inches(0.4), size=16, color=c, align=PP_ALIGN.CENTER)

    # right: bridge + browser
    rect(slide, Inches(8.4), Inches(2.0), Inches(3.0), Inches(1.3), C_PANEL)
    rect(slide, Inches(8.4), Inches(2.0), Inches(0.1), Inches(1.3), C_GREEN)
    tx(slide, "Web Bridge\nbridge.py",
       Inches(8.55), Inches(2.08), Inches(2.7), Inches(0.5),
       size=13, bold=True, color=C_GREEN)
    tx(slide, "Python asyncio\nUnix socket ↔ WebSocket",
       Inches(8.55), Inches(2.6), Inches(2.7), Inches(0.55),
       size=11, color=C_MUTED)

    rect(slide, Inches(8.4), Inches(3.8), Inches(3.0), Inches(1.3), C_PANEL)
    rect(slide, Inches(8.4), Inches(3.8), Inches(0.1), Inches(1.3), C_CYAN)
    tx(slide, "ブラウザ HW パネル",
       Inches(8.55), Inches(3.88), Inches(2.7), Inches(0.5),
       size=13, bold=True, color=C_CYAN)
    tx(slide, "LED / Button / OLED\nをブラウザで操作",
       Inches(8.55), Inches(4.4), Inches(2.7), Inches(0.55),
       size=11, color=C_MUTED)

    # arrows
    tx(slide, "←→", Inches(8.1), Inches(2.55), Inches(0.45), Inches(0.4),
       size=14, color=C_GREEN, align=PP_ALIGN.CENTER)
    tx(slide, "↕", Inches(9.7), Inches(3.35), Inches(0.4), Inches(0.4),
       size=16, color=C_GREEN, align=PP_ALIGN.CENTER)

    # bottom: RasPi5
    rect(slide, Inches(3.5), Inches(6.3), Inches(6.0), Inches(0.95), C_BG)
    rect(slide, Inches(3.5), Inches(6.3), Inches(6.0), Inches(0.1), C_GREEN)
    tx(slide, "Raspberry Pi 5  —  同一バイナリ、shim なし、実デバイスで動作",
       Inches(3.6), Inches(6.45), Inches(5.8), Inches(0.7),
       size=14, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    tx(slide, "↓  scp", Inches(6.2), Inches(5.8), Inches(1.0), Inches(0.4),
       size=13, color=C_GREEN, align=PP_ALIGN.CENTER)


# ── Slide 7: 検証結果 ─────────────────────────────────────────────────────
def s_results(prs):
    slide = blank(prs)
    bg(slide)
    top_bar(slide, C_GREEN)

    label(slide, "06  実証結果", C_GREEN)
    heading(slide, "動いた。コード変更ゼロで実機に展開できた",
            "PoC として想定した全シナリオの検証が完了")

    rows = [
        ["GPIO  LED + ボタン",    "✅ シム動作",  "✅ 実機動作", "完了"],
        ["I2C   VL53L0X センサー", "✅ レジスタ模倣", "✅ 実デバイス", "完了"],
        ["SPI   MFRC-522 RFID",  "✅ レジスタ模倣", "✅ カード読取", "完了"],
        ["I2C   SSD1306 OLED",   "✅ Canvas 表示", "⏳ モジュール待ち", "進行中"],
        ["WebSocket HW パネル",   "✅ リアルタイム", "—", "完了"],
    ]
    headers = ["コンポーネント", "EC2 シミュレーション", "RasPi5 実機", "状態"]
    col_x = [Inches(0.5), Inches(5.2), Inches(8.0), Inches(11.0)]
    col_w = [Inches(4.5), Inches(2.6), Inches(2.8), Inches(2.0)]

    # header
    rect(slide, Inches(0.5), Inches(2.1), Inches(12.3), Inches(0.55), C_PANEL)
    for j, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
        tx(slide, h, x + Inches(0.1), Inches(2.17), w, Inches(0.4),
           size=12, bold=True, color=C_CYAN)

    for i, row in enumerate(rows):
        y = Inches(2.65 + i * 0.78)
        bg_c = C_PANEL if i % 2 == 0 else RGBColor(0x12, 0x1E, 0x36)
        rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.75), bg_c)
        for j, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
            c = (C_GREEN  if cell.startswith("✅") else
                 C_YELLOW if cell.startswith("⏳") else
                 C_GREEN  if cell == "完了" else
                 C_YELLOW if cell == "進行中" else C_WHITE)
            tx(slide, cell, x + Inches(0.1), y + Inches(0.18),
               w, Inches(0.45), size=12.5, color=c)

    # KPIs
    kpis = [
        (C_CYAN,   "3 種",    "HW インターフェース\n(GPIO / I2C / SPI)"),
        (C_GREEN,  "0 行",    "実機向けコード変更"),
        (C_PURPLE, "≈ $3/日", "EC2 シミュ稼働コスト"),
    ]
    for i, (c, num, lbl) in enumerate(kpis):
        x = Inches(0.5 + i * 4.1)
        y = Inches(6.7)
        rect(slide, x, y, Inches(3.85), Inches(0.68), C_PANEL)
        rect(slide, x, y, Inches(0.1),  Inches(0.68), c)
        tx(slide, num, x + Inches(0.25), y + Inches(0.05),
           Inches(1.3), Inches(0.55), size=24, bold=True, color=c)
        tx(slide, lbl, x + Inches(1.6), y + Inches(0.1),
           Inches(2.1), Inches(0.5), size=11, color=C_MUTED)


# ── Slide 8: AI × 組み込み ────────────────────────────────────────────────
def s_ai(prs):
    slide = blank(prs)
    bg(slide)
    top_bar(slide, C_YELLOW)

    label(slide, "07  なぜ今これが重要か", C_YELLOW)
    heading(slide, "AI が組み込み開発に入ってきた",
            "AI はクラウドにしか存在できない → クラウド上の HW 環境が必須になる")

    # Big quote box
    rect(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.6), C_PANEL)
    rect(slide, Inches(0.5), Inches(2.0), Inches(0.15), Inches(1.6), C_YELLOW)
    tx(slide,
       "Claude Code は今まさにこのプロジェクトのコードを書いている。\n"
       "だが Claude は実機の GPIO に触れない。\n"
       "→  仮想 HW 環境がなければ、AI はドライバを書いても動作確認できない。",
       Inches(0.8), Inches(2.1), Inches(11.8), Inches(1.4),
       size=16, italic=True, color=C_WHITE)

    use_cases = [
        (C_CYAN,   "ドライバ自動生成",
         "AI が I2C / SPI ドライバを生成\n→ CUSE スタブで即テスト\n実機なしで動作確認完結"),
        (C_PURPLE, "自動テスト生成",
         "センサー値を自動注入して\n境界値・異常値テスト\nCI/CD に組み込み可能"),
        (C_GREEN,  "トレース解析",
         "HW トレースログを AI に渡し\nバグ原因を自動推定\nデバッグ工数を大幅削減"),
        (C_ORANGE, "コードレビュー",
         "PR 時にドライバを AI がレビュー\nタイミング・メモリ安全性を検査\nHW 起因のバグを事前防止"),
    ]
    for i, (c, title, body) in enumerate(use_cases):
        x = Inches(0.5 + i * 3.1)
        y = Inches(3.85)
        rect(slide, x, y, Inches(2.9), Inches(3.3), C_PANEL)
        rect(slide, x, y, Inches(2.9), Inches(0.1), c)
        tx(slide, title, x + Inches(0.15), y + Inches(0.2),
           Inches(2.6), Inches(0.5), size=14, bold=True, color=c)
        tx(slide, body, x + Inches(0.15), y + Inches(0.8),
           Inches(2.6), Inches(2.3), size=12.5, color=C_MUTED)

    tx(slide,
       "→  「AI 対応」は後付けの機能ではなく、本アーキテクチャの設計思想そのもの",
       Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
       size=13, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)


# ── Slide 9: ロードマップ ──────────────────────────────────────────────────
def s_roadmap(prs):
    slide = blank(prs)
    bg(slide)
    top_bar(slide, C_CYAN)

    label(slide, "08  ロードマップ", C_CYAN)
    heading(slide, "PoC から製品へ",
            "実証済みのコアを基盤に段階的に拡張する")

    phases = [
        ("今  ✅", "Phase 1 — PoC 完成", C_GREEN, [
            "GPIO / I2C / SPI シム確立",
            "EC2 Graviton + ブラウザ HW パネル",
            "RasPi5 実機デプロイ検証済",
        ]),
        ("次  🔧", "Phase 2 — 機能拡張", C_CYAN, [
            "OLED シム / LCD HAT シム追加",
            "GitHub Actions CI/CD 統合",
            "複数ボードへの対応",
        ]),
        ("近  🤖", "Phase 3 — AI 統合", C_PURPLE, [
            "Claude Code 自動生成フロー",
            "AI テストシナリオ生成",
            "HW トレース自動解析",
        ]),
        ("将来  🌐", "Phase 4 — SaaS 化", C_YELLOW, [
            "仮想 HW 環境プラットフォーム",
            "企業向け組み込み CI/CD",
            "Edge AI 開発の標準環境へ",
        ]),
    ]

    for i, (when, title, c, items) in enumerate(phases):
        x = Inches(0.5 + i * 3.1)
        y = Inches(2.0)
        rect(slide, x, y, Inches(2.9), Inches(5.1), C_PANEL)
        rect(slide, x, y, Inches(0.13), Inches(5.1), c)

        tx(slide, when, x + Inches(0.2), y + Inches(0.12),
           Inches(2.5), Inches(0.4), size=12, bold=True, color=c)
        tx(slide, title, x + Inches(0.2), y + Inches(0.55),
           Inches(2.5), Inches(0.6), size=13, bold=True, color=C_WHITE)

        divider(slide, y + Inches(1.2), C_BG)

        for j, item in enumerate(items):
            tx(slide, f"•  {item}",
               x + Inches(0.2), y + Inches(1.35 + j * 1.15),
               Inches(2.5), Inches(1.0), size=12, color=C_MUTED)

    # timeline
    rect(slide, Inches(0.5), Inches(7.2), Inches(12.3), Inches(0.1), C_PANEL)
    for i, c in enumerate([C_GREEN, C_CYAN, C_PURPLE, C_YELLOW]):
        rect(slide, Inches(0.5 + i * 3.1), Inches(7.2), Inches(2.9), Inches(0.1), c)


# ── Slide 10: まとめ ──────────────────────────────────────────────────────
def s_summary(prs):
    slide = blank(prs)
    bg(slide)
    rect(slide, 0, 0, W, Inches(3.0), C_PANEL)
    top_bar(slide, C_CYAN)

    label(slide, "09  まとめ", C_CYAN)
    tx(slide, "費用対効果が、ついて逆転した",
       Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.9),
       size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    tx(slide,
       "ISA 統一・Cloud IDE 無料化・AI の台頭。3つが揃い「仮想HW開発環境」は今こそやるべき投資になった。",
       Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.7),
       size=15, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    vals = [
        (C_CYAN,   "開発速度",   "ハード依存ゼロ\n即日スタート可能"),
        (C_GREEN,  "品質",       "CI/CD で自動テスト\n実機と同等の検証"),
        (C_PURPLE, "コスト",     "テスト機材費ゼロ\n≈ $3/日で代替"),
        (C_YELLOW, "AI 対応",    "Claude Code と\nシームレスに統合"),
    ]
    for i, (c, title, body) in enumerate(vals):
        x = Inches(0.5 + i * 3.1)
        y = Inches(3.3)
        rect(slide, x, y, Inches(2.9), Inches(2.8), C_BG)
        rect(slide, x, y, Inches(2.9), Inches(0.1), c)
        tx(slide, title, x + Inches(0.15), y + Inches(0.2),
           Inches(2.6), Inches(0.5), size=18, bold=True, color=c,
           align=PP_ALIGN.CENTER)
        tx(slide, body, x + Inches(0.15), y + Inches(0.85),
           Inches(2.6), Inches(1.7), size=13, color=C_MUTED,
           align=PP_ALIGN.CENTER)

    rect(slide, Inches(1.5), Inches(6.45), Inches(10.3), Inches(0.78), C_PANEL)
    rect(slide, Inches(1.5), Inches(6.45), Inches(0.13), Inches(0.78), C_CYAN)
    tx(slide,
       "本 PoC をベースに、組み込み開発の新しいスタンダードを共に構築しましょう",
       Inches(1.8), Inches(6.55), Inches(9.8), Inches(0.55),
       size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    tx(slide, "ExperimentalDevEnv  ／  github.com/thousandsofties/experimentaldevenv",
       Inches(0.5), Inches(7.2), Inches(12.3), Inches(0.25),
       size=10, color=C_MUTED, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────
def main():
    prs = prs_new()
    s_title(prs)
    s_not_new(prs)
    s_why_not(prs)
    s_why_now(prs)
    s_solution(prs)
    s_arch(prs)
    s_results(prs)
    s_ai(prs)
    s_roadmap(prs)
    s_summary(prs)

    out = "/home/user/ExperimentalDevEnv/docs/ExperimentalDevEnv_Intro.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
